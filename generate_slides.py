#!/usr/bin/env python3
"""Generate PowerPoint slides for 'How Does Town Meeting Work?' video.

Uses the formatting of Town Meeting Member Orientation.pptx as a template:
- 10" x 5.625" (16:9) slide dimensions
- Town seal image in top right corner of every slide
- Clean title + body layout matching the orientation style
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os, io

DIR = os.path.dirname(__file__)
OUTPUT = os.path.join(DIR, "Town_Meeting_Video_Slides.pptx")
TEMPLATE = os.path.join(DIR, "Town Meeting Member Orientation.pptx")
LEAF_BLOWER_IMG = os.path.join(DIR, "gas_leaf_blower.jpg")
TURF_IMG = os.path.join(DIR, "artificial_turf_field.jpg")
TM_IMG = os.path.join(DIR, "town_meeting_2026.jpg")

# Colors matching the orientation style
TITLE_COLOR = RGBColor(0x1B, 0x3A, 0x5C)   # dark navy
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)      # off-white
BODY_BG = RGBColor(0xE8, 0xED, 0xF2)        # slightly darker for body
ACCENT = RGBColor(0x2E, 0x86, 0xAB)         # teal

# ---- Load template image from orientation file ----
def load_template_image():
    prs_template = Presentation(TEMPLATE)
    for master in prs_template.slide_masters:
        for shape in master.shapes:
            if hasattr(shape, 'image'):
                return shape.image.blob, shape.image.content_type
    raise RuntimeError("No image found in template master slide")

TEMPLATE_IMAGE_BYTES, TEMPLATE_IMAGE_MIME = load_template_image()
IMAGE_EXT = os.path.splitext(TEMPLATE_IMAGE_MIME.split('/')[-1])[1] or '.png'

# ---- Helpers ----

def add_logo(slide):
    """Add the town seal to the top-right corner of a slide."""
    with io.BytesIO(TEMPLATE_IMAGE_BYTES) as bio:
        slide.shapes.add_picture(bio, Inches(8.889), Inches(0), Inches(1.111), Inches(1.111))

def set_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=0.242):
    add_text_box(slide, 0.341, top, 8.3, 0.626, text,
                 font_size=28, bold=True, color=TITLE_COLOR)

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=20, color=DARK, line_spacing=1.8):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\u2022 " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_split(slide, left_title, left_items, right_title, right_items):
    """Two-column layout matching orientation style."""
    col_w = 4.0
    col_l = 0.8
    col_r = 5.3
    body_top = 1.2
    add_text_box(slide, col_l, body_top, col_w, 0.5, left_title,
                 font_size=24, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_bullets(slide, col_l, body_top + 0.6, col_w, 3.2, left_items, font_size=20)
    add_text_box(slide, col_r, body_top, col_w, 0.5, right_title,
                 font_size=24, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_bullets(slide, col_r, body_top + 0.6, col_w, 3.2, right_items, font_size=20)

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_logo(slide)
    return slide


# ============================================================
# CREATE PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)


# ============================================================
# SLIDE 1 — Title Card
# ============================================================
slide = new_slide(prs)
# Clear any existing shapes
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)
# Re-add logo
add_logo(slide)
set_bg(slide, WHITE)
add_text_box(slide, 0.8, 1.2, 8.4, 1.5,
             "How Does Town Meeting Work?\nTwo Big Decisions in Natick",
             font_size=36, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — What Is Town Meeting?
# ============================================================
slide = new_slide(prs)
add_title(slide, "What Is Town Meeting?")
add_bullets(slide, 0.341, 1.1, 8.5, 2.0, [
    "Elected members vote on big decisions for the town.",
    "Any resident can come and share their opinion.",
    "Towns vote on money, rules, and local projects.",
    "It\u2019s like a student council \u2014 for the whole town!",
])
# Flow diagram using shapes
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
box_w, box_h = 1.6, 0.7
gap = 0.55
total = 4 * box_w + 3 * gap
left_start = (10 - total) / 2
labels = ["Residents", "Town Meeting", "Vote", "Decision"]
y = 3.8
from lxml import etree
for i, label in enumerate(labels):
    x = left_start + i * (box_w + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = TITLE_COLOR
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    # Arrow between boxes
    if i < len(labels) - 1:
        ax = x + box_w + 0.05
        ay = y + box_h / 2
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(ax), Inches(ay), Inches(ax + gap - 0.1), Inches(ay))
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(2)
        # Add arrowhead via XML
        ln = conn.line._ln
        tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tailEnd)

# ============================================================
# SLIDE 3 — The Town Moderator
# ============================================================
slide = new_slide(prs)
add_title(slide, "The Town Moderator")
if os.path.exists(TM_IMG):
    slide.shapes.add_picture(TM_IMG, Inches(5.5), Inches(1.2), Inches(4.2), Inches(3.15))
add_bullets(slide, 0.341, 1.1, 5.0, 3.5, [
    "The Moderator keeps Town Meeting fair and organized.",
    "In Natick, that\u2019s me \u2014 Jeff Alderson!",
    "Like a referee: I keep order and give everyone a turn to speak.",
    "I explain each question and call the vote.",
    "I don\u2019t take sides \u2014 I stay neutral.",
])


# ============================================================
# SLIDE 4 — Two Big Questions
# ============================================================
slide = new_slide(prs)
add_title(slide, "Two Big Questions Before Town Meeting")
add_text_box(slide, 0.8, 1.3, 3.8, 0.5, "Gas Leaf Blowers",
             font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_bullets(slide, 0.8, 1.9, 3.8, 1.5, [
    "Should Natick phase them out?",
    "Debated in Fall 2024",
], font_size=20)
add_text_box(slide, 5.4, 1.3, 3.8, 0.5, "Artificial Turf",
             font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_bullets(slide, 5.4, 1.9, 3.8, 1.5, [
    "Should Natick pause new fields?",
    "Debated in Fall 2025",
], font_size=20)

# ============================================================
# SLIDE 5 — Gas Leaf Blowers
# ============================================================
slide = new_slide(prs)
add_title(slide, "Question 1: Gas Leaf Blowers")
# Place the leaf blower photo on the right side
if os.path.exists(LEAF_BLOWER_IMG):
    slide.shapes.add_picture(LEAF_BLOWER_IMG, Inches(5.5), Inches(1.2), Inches(4.2), Inches(3.15))
add_bullets(slide, 0.341, 1.1, 5.0, 3.5, [
    "EcoNatick asked Town Meeting to phase out gas leaf blowers.",
    "They are loud and pollute the air.",
    "They wanted people to switch to electric \u2014 or a rake!",
])

# ============================================================
# SLIDE 6 — Leaf Blowers: The Debate
# ============================================================
slide = new_slide(prs)
add_title(slide, "The Debate: Leaf Blowers")
add_split(slide,
    "FOR a Phase-Out", [
        "Better for health \u2014 less pollution",
        "Quieter neighborhoods",
        "Electric blowers are getting better",
    ],
    "AGAINST a Phase-Out", [
        "Battery blowers not strong enough yet",
        "Too expensive for small businesses",
        "Some called it \u2018government overreach\u2019",
    ])

# ============================================================
# SLIDE 7 — Leaf Blowers: The Outcome
# ============================================================
slide = new_slide(prs)
add_title(slide, "Outcome: Sent for Study")
add_text_box(slide, 1.5, 1.4, 7.0, 1.0,
             "79  \u2014  27  \u2014  1",
             font_size=54, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.5, 7.0, 0.5,
             "Yes        No        Abstain",
             font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_bullets(slide, 0.8, 3.3, 8.0, 1.5, [
    "Town Meeting sent it to a committee to study more.",
    "No decision yet \u2014 they decided to learn more first.",
], font_size=22)

# ============================================================
# SLIDE 8 — Artificial Turf
# ============================================================
slide = new_slide(prs)
add_title(slide, "Question 2: Artificial Turf")
if os.path.exists(TURF_IMG):
    slide.shapes.add_picture(TURF_IMG, Inches(5.5), Inches(1.2), Inches(4.2), Inches(3.15))
add_bullets(slide, 0.341, 1.1, 5.0, 3.5, [
    "Some residents asked for a three-year pause on new turf fields.",
    "Artificial turf is plastic grass for sports fields.",
    "They wanted to study health and environmental risks first.",
])

# ============================================================
# SLIDE 9 — Artificial Turf: The Debate
# ============================================================
slide = new_slide(prs)
add_title(slide, "The Debate: Artificial Turf")
add_split(slide,
    "FOR a Pause", [
        "Microplastics may be bad for health",
        "Turf fields get very hot in summer",
        "Hard to recycle old turf",
    ],
    "AGAINST a Pause", [
        "Not enough grass fields for kids to play",
        "Turf can be used more hours than grass",
        "A delay could push back important projects",
    ])

# ============================================================
# SLIDE 10 — Artificial Turf: The Outcome
# ============================================================
slide = new_slide(prs)
add_title(slide, "Outcome: Three-Year Pause Passed!")
add_text_box(slide, 1.5, 1.4, 7.0, 1.0,
             "74  \u2014  28  \u2014  3",
             font_size=54, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.5, 7.0, 0.5,
             "Yes        No        Abstain",
             font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_bullets(slide, 0.8, 3.3, 8.0, 1.5, [
    "No new artificial turf fields for three years.",
    "The town will study the issue first.",
], font_size=22)

# ============================================================
# SLIDE 11 — What Can We Learn?
# ============================================================
slide = new_slide(prs)
add_title(slide, "What Can We Learn?")
lessons = [
    ("1. Anyone can speak up",
     "Bring your ideas \u2014 even if you\u2019re not a leader."),
    ("2. Debating helps",
     "Hearing both sides helps us make better choices."),
    ("3. Every vote counts",
     "Each person gets one vote. That\u2019s democracy!"),
]
y = 1.3
for title, body in lessons:
    add_text_box(slide, 0.8, y, 7.8, 0.5, title,
                 font_size=24, bold=True, color=ACCENT)
    add_text_box(slide, 0.8, y + 0.55, 7.8, 0.4, body,
                 font_size=20, color=DARK)
    y += 1.3

# Save
prs.save(OUTPUT)
print(f"Slides saved to: {OUTPUT}")
